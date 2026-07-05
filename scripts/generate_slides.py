#!/usr/bin/env python
"""Génère la présentation de soutenance (.pptx) du projet de qualité de l’air.

Diaporama 16:9 moderne. Chaque diapositive porte peu de texte mais contient des
NOTES DU PRÉSENTATEUR (script à lire), pour qu’un étudiant puisse présenter
facilement sans être expert. Couvre le contexte, la démarche et les compétences
C9 à C13.

Lancer :  ./.venv/bin/python scripts/generate_slides.py
Sortie :  docs/fr/PRESENTATION.pptx
"""
from __future__ import annotations

import json
import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
CAP = ROOT / "docs" / "fr" / "captures"
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "docs" / "fr" / "PRESENTATION.pptx"

INK = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREY = RGBColor(0x47, 0x55, 0x69)
PANEL = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Calibri"

metrics = json.loads((ROOT / "reports" / "metrics.json").read_text())
BEST = metrics["best_model"]
BEST_R2 = metrics["results"][BEST]["test_r2"]


def png_size(path: Path) -> tuple[int, int]:
    with open(path, "rb") as fh:
        head = fh.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return (1, 1)
    return struct.unpack(">II", head[16:24])


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    _fill(sh, color)
    return sh


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, line_spacing=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if not text:
            continue  # paragraphe espaceur : pas de run vide
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return box


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, idx, total=14):
    rect(slide, Inches(0.6), Inches(7.04), Inches(0.28), Inches(0.05), ACCENT)
    txt(slide, Inches(0.95), Inches(6.88), Inches(7), Inches(0.4),
        [("Qualité de l’air  ·  Prédiction de l’AQI", 10, MUTED, False)])
    txt(slide, Inches(11.2), Inches(6.88), Inches(1.5), Inches(0.4),
        [(f"{idx:02d} / {total:02d}", 10, MUTED, False)], align=PP_ALIGN.RIGHT)


def header(slide, kicker, title):
    rect(slide, Inches(0.6), Inches(0.62), Inches(0.18), Inches(0.6), ACCENT)
    txt(slide, Inches(0.95), Inches(0.55), Inches(11), Inches(0.4),
        [(kicker.upper(), 13, ACCENT, True)])
    txt(slide, Inches(0.92), Inches(0.95), Inches(11.8), Inches(0.95),
        [(title, 28, INK, True)])


def en_clair(slide, x, y, w, text):
    rect(slide, x, y, Inches(0.07), Inches(0.85), ACCENT)
    txt(slide, x + Inches(0.2), y - Inches(0.02), w - Inches(0.2), Inches(0.9),
        [("À RETENIR", 11, ACCENT, True), (text, 15, GREY, False)])


def image_card(slide, path: Path, x, y, max_w, max_h):
    if not path.exists():
        return
    rect(slide, x, y, max_w, max_h, PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    pad = Inches(0.18)
    iw, ih = max_w - 2 * pad, max_h - 2 * pad
    w, h = png_size(path)
    ratio = w / h
    if ratio > iw / ih:
        width, height = iw, int(iw / ratio)
    else:
        height, width = ih, int(ih * ratio)
    slide.shapes.add_picture(str(path), x + pad + (iw - width) // 2,
                             y + pad + (ih - height) // 2, width=width, height=height)


# --------------------------------------------------------------------------
def slide_cover(prs):
    s = new_slide(prs)
    rect(s, 0, 0, EMU_W, EMU_H, INK)
    rect(s, Inches(0.95), Inches(2.0), Inches(1.0), Inches(0.12), ACCENT)
    txt(s, Inches(0.95), Inches(1.25), Inches(11.5), Inches(0.5),
        [("PROJET DE FIN DE FORMATION", 15, RGBColor(0x93, 0xB4, 0xF5), True)])
    txt(s, Inches(0.9), Inches(2.35), Inches(11.6), Inches(2.4),
        [("Prédiction de la qualité de l’air", 40, WHITE, True),
         ("et mise en production du modèle", 40, WHITE, True)],
        line_spacing=1.05)
    txt(s, Inches(0.95), Inches(4.85), Inches(11.2), Inches(0.9),
        [("Solution d’intelligence artificielle de bout en bout : "
          "modèle, API REST, application, supervision et tests automatisés.",
          18, RGBColor(0xCB, 0xD5, 0xE1), False)], line_spacing=1.1)
    txt(s, Inches(0.95), Inches(6.35), Inches(11.4), Inches(0.7),
        [("Compétences C9 à C13   ·   Titre Développeur en IA (RNCP)   ·   "
          "Marimoun   ·   2026", 13, MUTED, False)])
    notes(s,
          "Bonjour, je vais vous présenter mon projet de fin de formation. "
          "L’objectif est double : d’abord prédire la qualité de l’air "
          "à partir de capteurs, et surtout rendre ce modèle réellement "
          "utilisable, comme dans une entreprise. Je vais vous montrer les cinq "
          "grandes étapes, qui correspondent aux compétences C9 à C13.")


def slide_statement(prs, idx, kicker, big_lines, sub, note):
    s = new_slide(prs)
    rect(s, Inches(0.6), Inches(0.62), Inches(0.18), Inches(0.6), ACCENT)
    txt(s, Inches(0.95), Inches(0.55), Inches(11), Inches(0.4),
        [(kicker.upper(), 13, ACCENT, True)])
    txt(s, Inches(0.9), Inches(2.3), Inches(11.8), Inches(2.4),
        [(line, 32, INK, True) for line in big_lines], line_spacing=1.05)
    txt(s, Inches(0.95), Inches(4.9), Inches(11.2), Inches(1.4),
        [(sub, 18, GREY, False)], line_spacing=1.1)
    footer(s, idx)
    notes(s, note)


def slide_solution(prs, idx):
    s = new_slide(prs)
    header(s, "Démarche · Solution", "Cible recalculée et variables du modèle")
    txt(s, Inches(0.9), Inches(1.8), Inches(6.7), Inches(0.95),
        [("Cible : indice EPA, fonction déterministe des concentrations de "
          "polluants (en remplacement de la colonne bruitée).", 15, INK, True)],
        line_spacing=1.05)
    txt(s, Inches(0.9), Inches(2.85), Inches(6.7), Inches(0.4),
        [("Variables explicatives : 30 caractéristiques", 15, ACCENT, True)])
    fams = ["Polluants (7) : CO, NOx, NO2, O3, SO2, PM2.5, PM10",
            "Météo (5) : température, humidité, pression, vent (vitesse, direction)",
            "Dérivées (6) : ratios CO/NOx et NOx/NO2, moyennes mobiles, indice T-H",
            "Temporelles : heure, jour, mois, saison + encodages cycliques sin/cos",
            "Interactions (2) : PM2.5 × Température, O3 × Humidité"]
    box = s.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(6.7), Inches(3.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, fam in enumerate(fams):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = "•  " + fam
        r.font.size = Pt(13)
        r.font.name = FONT
        r.font.color.rgb = GREY
    image_card(s, FIG / "feature_importance.png",
               Inches(7.8), Inches(1.95), Inches(4.9), Inches(4.5))
    txt(s, Inches(7.8), Inches(6.5), Inches(4.9), Inches(0.3),
        [("Importance des variables : O3, PM2.5 et PM10 dominent", 11, MUTED, False)],
        align=PP_ALIGN.CENTER)
    footer(s, idx)
    notes(s,
          "[~2:00] La solution a été de recalculer un vrai indice de qualité de "
          "l’air avec la formule officielle de l’EPA, à partir des polluants. Le "
          "modèle s’appuie sur trente variables, regroupées en familles : les sept "
          "polluants, cinq variables météo, six variables dérivées (ratios, "
          "moyennes mobiles, indice température-humidité), des variables "
          "temporelles avec leurs encodages cycliques, et deux interactions "
          "physiques. Le graphique de droite confirme que les variables les plus "
          "importantes sont l’ozone et les particules fines, ce qui correspond à "
          "la logique de l’AQI.")


def slide_chiffres(prs, idx):
    s = new_slide(prs)
    header(s, "Résultats", "Trois indicateurs clés")
    stats = [(f"{BEST_R2:.3f}".replace(".", ","), "R² (jeu de test)",
              "ajustement quasi parfait"),
             ("33", "tests automatisés", "rejoués en intégration continue"),
             ("C9–C13", "compétences couvertes", "de l’API à la mise en production")]
    x = Inches(0.9)
    for value, label, hint in stats:
        rect(s, x, Inches(2.3), Inches(3.7), Inches(3.0), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, Inches(2.6), Inches(3.7), Inches(1.3),
            [(value, 46, ACCENT, True)], align=PP_ALIGN.CENTER)
        txt(s, x, Inches(3.95), Inches(3.7), Inches(0.5),
            [(label, 16, INK, True)], align=PP_ALIGN.CENTER)
        txt(s, x, Inches(4.5), Inches(3.7), Inches(0.7),
            [(hint, 13, GREY, False)], align=PP_ALIGN.CENTER)
        x += Inches(4.0)
    footer(s, idx)
    notes(s,
          "Voici les résultats résumés en trois chiffres simples. Premièrement, "
          "le modèle obtient un score R² de 0,999 ; un score proche de 1 veut "
          "dire qu’il se trompe très peu. Deuxièmement, j’ai 33 tests "
          "automatiques qui vérifient que tout fonctionne. Et troisièmement, les "
          "cinq compétences attendues, de C9 à C13, sont toutes couvertes.")


def slide_comp(prs, idx, kicker, title, lines, clair, image, note):
    s = new_slide(prs)
    header(s, kicker, title)
    runs = [(lines[0], 18, INK, True)] + [(ln, 16, GREY, False) for ln in lines[1:]]
    txt(s, Inches(0.95), Inches(2.0), Inches(5.9), Inches(2.8), runs, line_spacing=1.08)
    en_clair(s, Inches(0.95), Inches(5.2), Inches(5.7), clair)
    image_card(s, image, Inches(7.1), Inches(1.95), Inches(5.6), Inches(4.6))
    footer(s, idx)
    notes(s, note)


def slide_pipeline(prs, idx):
    s = new_slide(prs)
    header(s, "C13 · MLOps", "Tout est automatisé (approche MLOps)")
    txt(s, Inches(0.95), Inches(2.0), Inches(11.4), Inches(0.7),
        [("Chaîne déclenchée automatiquement à chaque modification du code :",
          18, INK, True)])
    steps = ["Vérifier\nle code", "Lancer\nles tests", "Réentraîner\nle modèle",
             "Construire\nl’image", "Publier"]
    x = Inches(0.9)
    for i, st in enumerate(steps):
        rect(s, x, Inches(3.1), Inches(2.0), Inches(1.25), PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, Inches(3.3), Inches(2.0), Inches(0.95),
            [(st, 14, INK, True)], align=PP_ALIGN.CENTER, line_spacing=1.0)
        if i < len(steps) - 1:
            txt(s, x + Inches(1.98), Inches(3.35), Inches(0.45), Inches(0.8),
                [("→", 22, ACCENT, True)], align=PP_ALIGN.CENTER)
        x += Inches(2.42)
    en_clair(s, Inches(0.95), Inches(5.2), Inches(11.4),
             "Aucune intervention manuelle : reproductibilité et traçabilité de bout en bout.")
    footer(s, idx)
    notes(s,
          "[~1:30] Dernière compétence : l’automatisation, le MLOps. À chaque "
          "modification du code, une chaîne se déclenche seule : elle vérifie le "
          "style, lance les 33 tests, réentraîne le modèle, construit le conteneur "
          "Docker et le publie. Le pipeline DVC rend tout reproductible. Libellé "
          "officiel C13 : créer une chaîne de livraison continue d’un modèle d’IA "
          "dans une approche MLOps, pour automatiser les étapes de validation, de "
          "test, de packaging et de déploiement du modèle.")


def slide_merci(prs):
    s = new_slide(prs)
    rect(s, 0, 0, EMU_W, EMU_H, INK)
    rect(s, Inches(0.95), Inches(3.5), Inches(1.0), Inches(0.12), ACCENT)
    txt(s, Inches(0.95), Inches(2.55), Inches(11.5), Inches(1.0),
        [("Merci de votre attention", 42, WHITE, True)])
    txt(s, Inches(0.97), Inches(3.85), Inches(11), Inches(0.6),
        [("Place à la démonstration et à vos questions", 18, RGBColor(0xCB, 0xD5, 0xE1), False)])
    notes(s,
          "Voilà, je vous remercie de votre attention. Pour résumer : je n’ai "
          "pas seulement entraîné un modèle, j’ai construit toute la chaîne qui "
          "le rend fiable et utilisable. Je suis maintenant disponible pour une "
          "démonstration en direct et pour répondre à vos questions.")


def slide_agenda(prs, idx):
    s = new_slide(prs)
    header(s, "Plan", "Déroulé de la présentation (≈ 20 min)")
    items = [("Contexte et données", "2 min"),
             ("Solution et résultats", "3 min"),
             ("C9 · API REST", "2 min"),
             ("C10 · Application", "2 min"),
             ("C11 · Monitoring", "2 min"),
             ("C12 · Tests automatisés", "2 min"),
             ("C13 · Chaîne MLOps", "2 min"),
             ("Démonstration (vidéo)", "3 min"),
             ("Conclusion et questions", "2 min")]
    box = s.shapes.add_textbox(Inches(0.95), Inches(1.85), Inches(11.4), Inches(4.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (label, dur) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {label}"
        r.font.size = Pt(18)
        r.font.name = FONT
        r.font.color.rgb = INK
        r2 = p.add_run()
        r2.text = f"      ({dur})"
        r2.font.size = Pt(14)
        r2.font.name = FONT
        r2.font.color.rgb = ACCENT
    footer(s, idx)
    notes(s,
          "[~0:45] Voici le déroulé. Je commence par le contexte et la démarche, "
          "puis je détaille les cinq compétences C9 à C13, je vous montre une "
          "démonstration en vidéo, et je termine par la conclusion et vos "
          "questions. L’ensemble tient en une vingtaine de minutes.")


def slide_contexte(prs, idx):
    s = new_slide(prs)
    header(s, "Contexte", "Problématique et données")
    txt(s, Inches(0.95), Inches(1.95), Inches(7.0), Inches(3.4),
        [("Objectif : estimer l’indice de qualité de l’air (AQI) à partir de "
          "capteurs.", 18, INK, True),
         ("Jeu de données : 4 000 observations horaires (polluants et météo).",
          16, GREY, False),
         ("", 6, GREY, False),
         ("Constat : la cible fournie est décorrélée des variables (|r| < 0,04).",
          18, INK, True),
         ("Conséquence : R² négatif pour tous les modèles testés.",
          16, GREY, False)],
        line_spacing=1.06)
    rect(s, Inches(8.5), Inches(2.0), Inches(4.0), Inches(3.0), PANEL,
         MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(8.5), Inches(2.45), Inches(4.0), Inches(1.5),
        [("< 0,04", 54, ACCENT, True)], align=PP_ALIGN.CENTER)
    txt(s, Inches(8.5), Inches(3.95), Inches(4.0), Inches(0.9),
        [("corrélation maximale\ncible / variables", 14, GREY, False)],
        align=PP_ALIGN.CENTER)
    en_clair(s, Inches(0.95), Inches(5.7), Inches(11.4),
             "Une cible décorrélée des variables est, par construction, non apprenable.")
    footer(s, idx)
    notes(s,
          "[~2:00] Le but est d’estimer la qualité de l’air à partir de capteurs. "
          "Le jeu de données contient 4 000 mesures horaires. En l’explorant, j’ai "
          "découvert un problème : la colonne à prédire était du bruit, sans aucun "
          "lien avec les mesures (corrélation inférieure à 0,04). Résultat, tous "
          "les modèles échouaient. La diapo suivante explique comment je l’ai "
          "corrigé.")


def slide_architecture(prs, idx):
    s = new_slide(prs)
    header(s, "Architecture", "De la donnée au service en production")
    boxes = ["Données\n+ Modèle", "API REST\n(FastAPI)", "Application\n(Streamlit)",
             "Monitoring\n(métriques + alerte)", "MLOps\n(DVC, CI/CD)"]
    x = Inches(0.9)
    for i, b in enumerate(boxes):
        rect(s, x, Inches(2.7), Inches(2.0), Inches(1.5), PANEL,
             MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, Inches(2.95), Inches(2.0), Inches(1.1),
            [(b, 13, INK, True)], align=PP_ALIGN.CENTER, line_spacing=1.0)
        if i < len(boxes) - 1:
            txt(s, x + Inches(1.98), Inches(3.05), Inches(0.45), Inches(0.8),
                [("→", 22, ACCENT, True)], align=PP_ALIGN.CENTER)
        x += Inches(2.42)
    en_clair(s, Inches(0.95), Inches(5.0), Inches(11.4),
             "Chemin d’inférence unique et partagé (inference.py) : aucune "
             "divergence possible entre l’API et l’application.")
    footer(s, idx)
    notes(s,
          "[~1:30] Voici l’architecture d’ensemble. À gauche, la donnée et le "
          "modèle. Le modèle est exposé par une API, consommée par une "
          "application. Le tout est surveillé par le monitoring et industrialisé "
          "par la chaîne MLOps. Point clé : l’API et l’application partagent le "
          "même code de préparation des données, donc elles ne peuvent pas donner "
          "des résultats différents.")


def slide_demo(prs, idx):
    s = new_slide(prs)
    header(s, "Démonstration", "La solution en fonctionnement")
    txt(s, Inches(0.9), Inches(1.95), Inches(4.2), Inches(4.6),
        [("Contenu de la démonstration :", 16, INK, True),
         ("", 8, GREY, False),
         ("1. API : documentation /docs et appel /predict.", 15, GREY, False),
         ("", 6, GREY, False),
         ("2. Application : saisie, AQI et catégorie.", 15, GREY, False),
         ("", 6, GREY, False),
         ("3. Monitoring : /metrics et alerte de dérive.", 15, GREY, False)],
        line_spacing=1.05)
    px, py, pw = Inches(5.4), Inches(1.95), Inches(7.3)
    ph = Inches(7.3 * 9 / 16)
    rect(s, px, py, pw, ph, RGBColor(0x1E, 0x29, 0x3B), MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, px, py + Inches(0.75), pw, Inches(1.4), [("▶", 54, WHITE, True)],
        align=PP_ALIGN.CENTER)
    txt(s, px, py + Inches(2.2), pw, Inches(0.6),
        [("Insérez ici la vidéo de démonstration", 16, WHITE, True)],
        align=PP_ALIGN.CENTER)
    txt(s, px, py + Inches(2.85), pw, Inches(0.6),
        [("PowerPoint : Insertion › Vidéo › Cet appareil", 12,
          RGBColor(0xCB, 0xD5, 0xE1), False)], align=PP_ALIGN.CENTER)
    footer(s, idx)
    notes(s,
          "[~3:00] Je lance maintenant la vidéo de démonstration. Pendant qu’elle "
          "tourne, je commente : on voit d’abord la documentation de l’API et un "
          "appel qui renvoie l’indice ; ensuite l’application où l’on saisit des "
          "mesures et où s’affiche la catégorie ; enfin le monitoring avec ses "
          "métriques et une alerte de dérive déclenchée. Plan B si la vidéo ne se "
          "lance pas : je montre les mêmes étapes en direct ou via les captures "
          "des diapos précédentes.")


def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    slide_cover(prs)
    slide_agenda(prs, 2)
    slide_contexte(prs, 3)
    slide_solution(prs, 4)
    slide_chiffres(prs, 5)
    slide_architecture(prs, 6)
    slide_comp(prs, 7, "C9 · API REST", "Une API pour exposer le modèle",
               ["API REST développée avec FastAPI (serveur uvicorn).",
                "Validation des entrées, clé d’API, CORS, limitation de débit.",
                "Documentation OpenAPI générée ; interface standardisée inter-composants."],
               "Interface standardisée et sécurisée entre le modèle et le reste "
               "du système.",
               CAP / "c9_swagger_overview.png",
               "[~2:00] J’expose le modèle avec une API REST développée avec "
               "FastAPI. Elle valide les données reçues (et refuse une entrée hors "
               "plage), exige une clé d’accès, gère le CORS et limite le débit. Sa "
               "documentation OpenAPI est générée automatiquement, ce que vous "
               "voyez à l’écran. C’est elle qui permet aux autres composants de "
               "parler au modèle. Libellé officiel C9 : développer une API REST "
               "exposant un modèle d’IA en respectant ses spécifications "
               "fonctionnelles et techniques et les standards de qualité et de "
               "sécurité, pour permettre l’interaction entre le modèle et les "
               "autres composants du projet.")
    slide_comp(prs, 8, "C10 · Application", "Intégrer l’API dans une application",
               ["Application web (Streamlit) consommant l’API via le contrat OpenAPI.",
                "Saisie des mesures, restitution de l’AQI et de sa catégorie.",
                "Accessibilité (RGAA) : information par texte, icône et couleur."],
               "Accès métier au modèle, sans compétence technique requise.",
               CAP / "c10_result.png",
               "[~2:00] J’ai intégré l’API dans une application web Streamlit, en "
               "m’appuyant sur la documentation technique (le contrat OpenAPI). "
               "L’utilisateur saisit des mesures et obtient l’indice et sa "
               "catégorie, à l’unité ou par lot. J’ai respecté les normes "
               "d’accessibilité : la catégorie est restituée par texte, icône et "
               "couleur, jamais la couleur seule. Libellé officiel C10 : intégrer "
               "l’API d’un modèle d’IA dans une application, en respectant les "
               "spécifications et les normes d’accessibilité, à l’aide de la "
               "documentation technique de l’API.")
    slide_comp(prs, 9, "C11 · Monitoring", "Surveiller le modèle dans le temps",
               ["Collecte : métriques Prometheus en temps réel (/metrics).",
                "Restitution : rapport de dérive des données (Evidently).",
                "Alerte automatique sur dépassement de seuil (amélioration itérative)."],
               "Surveillance continue pour détecter et corriger la dégradation "
               "du modèle.",
               CAP / "c11_drift_report.png",
               "[~2:00] Je surveille le modèle en continu. Côté collecte, j’expose "
               "des métriques Prometheus. Côté restitution, je produis un rapport "
               "de dérive des données. Et côté alerte, dès que les données "
               "s’éloignent trop de l’entraînement, une alerte se déclenche "
               "automatiquement : c’est ce qui permet d’améliorer le modèle de "
               "façon itérative. Libellé officiel C11 : monitorer un modèle d’IA à "
               "partir de métriques courantes et spécifiques, en intégrant les "
               "outils de collecte, d’alerte et de restitution, pour permettre "
               "l’amélioration itérative du modèle.")
    slide_comp(prs, 10, "C12 · Tests", "Garantir la qualité par les tests",
               ["Règles de validation des jeux de données (pandera).",
                "Tests des étapes : préparation, entraînement, évaluation.",
                "Barrière de validation du modèle (R² > 0,9) pour la CI."],
               "Qualité garantie et automatisée avant toute mise en production.",
               FIG / "pred_vs_actual.png",
               "[~1:30] J’ai écrit 33 tests automatiques qui couvrent toute la "
               "chaîne : des règles de validation des données, les étapes de "
               "préparation, d’entraînement et d’évaluation, et une barrière de "
               "validation du modèle qui bloque si le R² descend sous 0,9. Tout "
               "est rejoué par l’intégration continue. Libellé officiel C12 : "
               "programmer les tests automatisés en définissant les règles de "
               "validation des jeux de données, des étapes de préparation, "
               "d’entraînement, d’évaluation et de validation du modèle, pour "
               "l’intégration continue et un niveau de qualité élevé.")
    slide_pipeline(prs, 11)
    slide_demo(prs, 12)
    slide_statement(prs, 13, "Conclusion",
                    ["Au-delà du modèle :",
                     "la chaîne complète de mise en production."],
                    "Perspectives : déploiement cloud, accessibilité renforcée "
                    "et suivi des performances en conditions réelles.",
                    "[~1:00] Pour conclure : la vraie valeur du projet, ce n’est "
                    "pas seulement le modèle, c’est toute la chaîne qui le rend "
                    "fiable et utilisable en production, des compétences C9 à C13. "
                    "Si je devais continuer, je déploierais la solution dans le "
                    "cloud et je suivrais ses performances réelles dans le temps.")
    slide_merci(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Présentation écrite : {OUT}  ({len(prs.slides._sldIdLst)} diapositives)")


if __name__ == "__main__":
    build()
